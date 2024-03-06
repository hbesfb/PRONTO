#!/usr/bin/env python
# -*- coding: utf-8 -*-
#
# This script needs to be run with python above version 3.
# To install module xlutils, run the command: sudo pip install xlutils
# To install module python-pptx, run the command: sudp pip3 install python-pptx
"""
This script is a tool used to filter and analysis data from TSO500 results.
And generate the PP report based on the results data and template file.
"""

import os
import re
import sys
import shutil
import getopt
import time
from configparser import ConfigParser
from xlutils.copy import copy
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE
from decimal import Decimal

runID = ""
DNA_sampleID = ""
RNA_sampleID = ""
extra_path = ""
batch_nr = ""
tumor_content_nr = ""
ipd_birth_year = ""
ipd_diagnosis_year = "-"
ipd_age = ""
ipd_gender = ""
ipd_consent = ""
ipd_collection_year = "-"
requisition_hospital = ""
ipd_material_id = ""
DNA_material_id = ""
RNA_material_id = ""
requisition_hospital = ""
extraction_hospital = ""
inclusion_site = ""
ipd_clinical_diagnosis = "-"
sample_material = ""
sample_type = ""
tumor_type = ""
TMB_DRUP = ""
str_TMB_DRUP = ""
TMB_TSO500 = ""
MSI_TSO500 = ""
pipline = ""

def read_exl(data_file,filter_column,key_word):
	data = []
	mark = False
	nrow_mark = 0
	col0 = "Sample_ID"
	global DNA_sampleID
	for line in open(data_file):
		line_cells = line.split('\t')
		if(line.split('\t')[0] == col0 and not mark):
			line_cells = line.split('\t')
			for col in range(len(line_cells)):
				if(line_cells[col] == "IGV_QC"):
					IGV_QC_col = col
				if(line_cells[col] == "Class_judgement"):
					Class_judgement_col = col
				if(line_cells[col] == filter_column):
					filter_column_n = col
			line_cells_string = [line_cells[i] + '\t' for i in range(len(line_cells))]
			line_cells_string.append('\n')
			data.append(line_cells_string)
			mark = True
		if(line.startswith(DNA_sampleID) and mark):
			if(line_cells[IGV_QC_col] == "Not OK" and line_cells[Class_judgement_col] != "exclude"):
				print ("""              Dataset error: 
			IGV_QC is 'Not OK', but Class_judement is not 'exclude'. Please check the QC Excel file and fix the mistake before run this script again!
				""")
				sys.exit(0)
			filter_column_data = line_cells[filter_column_n]
			if key_word.startswith('!'):
				if("&&" in key_word):
					not_key = key_word.replace('!','')
					key = not_key.split(" && ")
				else:
					key = key_word.split('!')
				appear = False
				for filter_column in filter_column_data.split(','):
					if(filter_column in key):
						appear = True
				if(appear == False):
					line_cells_string = [line_cells[i] + '\t' for i in range(len(line_cells))]
					line_cells_string.append('\n')
					data.append(line_cells_string)
			else:
				key = key_word.split(',')
				if(line_cells[filter_column_n] in key):
					line_cells_string = [line_cells[i] + '\t' for i in range(len(line_cells))]
					line_cells_string.append('\n')
					data.append(line_cells_string)
	return data


def read_exl_col(data_file,filter_column,key_word,columns,MTB_format):
	mark = False
	nrow_mark = 0
	col0 = "Sample_ID"
	global DNA_sampleID
	nlines = len(open(data_file).readlines())
	data = [[] for n in range(nlines)]
	column_mark = []
	columnNames = columns.split(',')
	d = 0
	for line in open(data_file):
		line_cells = line.split('\t')
		if(line_cells[0] == col0 and not mark):
			for col in range(len(line_cells)):
				if(line_cells[col] == "IGV_QC"):
					IGV_QC_col = col
				if(line_cells[col] == "Class_judgement"):
					Class_judgement_col = col
				if(line_cells[col] == filter_column):
					filter_column_n = col
				if(line_cells[col] == "Coding_status"):
					Coding_status_col = col
				if(line_cells[col] == "Genomic_location"):
					Genomic_location_col = col
				if(line_cells[col] == "DNA_change"):
					DNA_change_col = col
				for m in range(len(columnNames)):
					if(line_cells[col] in columnNames):
						line_cells_string = line_cells[col] + '\t'
						data[0].append(line_cells_string)
						column_mark.append(col)
						break
			if(MTB_format == True):
				data[0].insert(1,"MTB_formart\t")
			data[0].append('\n')
			mark = True
		if(line.startswith(DNA_sampleID) and mark):
			if(line_cells[IGV_QC_col] == "Not OK" and line_cells[Class_judgement_col] != "exclude"):
				print ("""              Dataset error: 
			IGV_QC is 'Not OK', but Class_judement is not 'exclude'. Please check the QC Excel file and fix the mistake before run this script again!
				 """)
				sys.exit(0)
			filter_column_data = line_cells[filter_column_n]
			if key_word.startswith('!'):
				if("&&" in key_word):
					not_key = key_word.replace('!','')
					key = not_key.split(" && ")
				else:
					key = key_word.split('!')
				appear = False
				for filter_column in filter_column_data.split(','):
					if(filter_column in key):
						appear = True
				if(appear == False):
					d += 1
					for num in column_mark:
						if(num == Coding_status_col):
							line_cells[num] = line_cells[num].replace("_variant", "") + '\t'
						else:
							line_cells[num] = line_cells[num] + '\t'
						data[d].append(line_cells[num])
					if(MTB_format == True):
						MTB_format_str = "chr" + line_cells[Genomic_location_col].split(":")[0] + ":g." + line_cells[Genomic_location_col].split(":")[1].replace('\t','') + line_cells[DNA_change_col] + '\t'
						data[d].insert(1,MTB_format_str)
			else:
				key = key_word.split(',')
				if(line_cells[filter_column_n] in key):
					d += 1
					for num in column_mark:
						if(num == Coding_status_col):
							line_cells[num] = line_cells[num].replace("_variant", "") + '\t'
						else:
							line_cells[num] = line_cells[num] + '\t'
						data[d].append(line_cells[num])
			data[d].append('\n')
			
	return data


def filter_depth_tumor_all_col(data_config,depth_tumor_DNA):
	data = []
	data.append(data_config[0])
	p = data_config[0].index('Depth_tumor_DNA\t')
	for row in data_config:
		if(row[p] != 'Depth_tumor_DNA\t' and row[p] != ''):
			num = row[p].split('\t')[0]
			if(int(num) >= depth_tumor_DNA):
				data.append(row)
	return data


def filter_depth_tumor_cols(data_config,depth_tumor_DNA):
	data = [[] for n in range(len(data_config))]
	data[0] = data_config[0]
	p = data_config[0].index('Depth_tumor_DNA\t')
	for i in range(len(data_config)):
		for j in range(len(data_config[i])):
			if(data_config[i][p] != 'Depth_tumor_DNA\t' and data_config[i][p] != ''):
				num = data_config[i][p].split('\t')[0]
				if(int(num) >= depth_tumor_DNA):
					data[i].append(data_config[i][j])
	return data


def write_exl(output_file,data):
	file_dir = os.path.split(output_file)[0]
	if not os.path.exists(file_dir):
		os.makedirs(file_dir)
	txt_file = open(output_file, mode='w', encoding='utf-8')
	for item in data:
		data_string = str(item) + "\t"
		txt_file.writelines(item)
	txt_file.close()


def clear_blank_line(file_in,file_out):
	fr = open(file_in, 'r')
	fw = open(file_out, 'w')
	for line in fr.readlines():
		if(line.split()):
			fw.write(line)
	fr.close()
	fw.close()
	os.remove(file_in)


def get_patient_info_from_MTF(ipd_material_file,ipd_no,DNA_sampleID,RNA_sampleID):
	import xlrd
	global ipd_birth_year
	global ipd_clinical_diagnosis
	global ipd_gender
	global ipd_consent
	global DNA_material_id
	global RNA_material_id
	global ipd_collection_year
	global requisition_hospital
	global extraction_hospital
	global batch_nr
	global tumor_content_nr
	global inclusion_site
	open_exl_material = xlrd.open_workbook(ipd_material_file)
	sheet_material = open_exl_material.sheet_by_index(0)
	nrows_material = sheet_material.nrows
	ncols_material = sheet_material.ncols
	col_ipd = "InPreD ID"
	col_gender = "Gender"
	col_age = "Age"
	col_birth_date = "Date of birth"
	col_requisition_hospital = "Requester Hospital"
	col_material_id = "Sample material ID"
	col_consent = "Study ID"
	col_tumor_content_nr = "Tumor cells [%]"
	col_sampleID_material_id = "Sample ID"
	col_ex_sample_info = "Sample information"
	col_ex_data_section = "Extraction Data"
	col_ex_library_pre = "Library Preparation (LP) Data"
	col_extraction_hospital = "Extraction Hospital"
	col_batch_nr = "LP batch"
	col_clinical_diagnosis = "Clinical diagnosis"
	ipd_birth_date = ""
	sample_info_row = 0
	extra_data_row = 0
	library_pre_row = 0
	for l in range(nrows_material):
		if(sheet_material.cell_value(l,0) == col_ex_sample_info):
			sample_info_row = l
		if(sheet_material.cell_value(l,0) == col_ex_data_section):
			extra_data_row = l
		if(sheet_material.cell_value(l,0) == col_ex_library_pre):
			library_pre_row = l  
	for r in range(nrows_material):
		for c in range(ncols_material):
			if(sheet_material.cell_value(r,c) == col_ipd):
				ipd_MTF = sheet_material.cell_value(r+2,c)
				if(ipd_MTF != ipd_no):
					print("""               Error:
                        The InPreD patient ID in IPD Material Transit Form InPreD NGS file does not match with the IPD number! 
                        Please check and fix the mistake before run this script again!""")
					print("                 IPD is " + ipd_MTF + " in MTF, while IPD is " + ipd_no[3:] + " in TSO500.")
					sys.exit(0)
			if(sheet_material.cell_value(r,c) == col_birth_date):
				ipd_birth_date_exl = sheet_material.cell_value(r+2,c)
				try:
					datetime_date = str(xlrd.xldate_as_datetime(ipd_birth_date_exl,0))
					ipd_birth_year = datetime_date.split('-')[0]
				except:
					ipd_birth_year = "-"
			if(sheet_material.cell_value(r,c) == col_gender and ipd_gender == ""):
				ipd_gender = str(sheet_material.cell_value(r+2,c))
			if(sheet_material.cell_value(r,c) == col_age):
				ipd_age = str(sheet_material.cell_value(r+2,c))
			if(sheet_material.cell_value(r,c) == col_clinical_diagnosis):
				ipd_clinical_diagnosis = str(sheet_material.cell_value(r+2,c))
			if(sheet_material.cell_value(r,c) == col_consent and ipd_consent == ""):
				ipd_consent = str(sheet_material.cell_value(r+2,c))
				for r in range(r,(sample_info_row-2)):
					if(ipd_consent == "0.0"):
						ipd_consent = ""
					if(sheet_material.cell_value(r,6) == col_requisition_hospital and requisition_hospital == ""):
						requisition_hospital = sheet_material.cell_value(r+2,6)
					if((sheet_material.cell_value(r+2,c) != "" or sheet_material.cell_value(r+2,c) != "-" or sheet_material.cell_value(r+2,c) != "0.0") and str(sheet_material.cell_value(r+2,c)) not in ipd_consent):
						if(ipd_consent == ""):
							ipd_consent = str(sheet_material.cell_value(r+2,c))
						else:
							ipd_consent = ipd_consent + "," + str(sheet_material.cell_value(r+2,c))
						continue
			if(sheet_material.cell_value(r,c) == col_material_id and ipd_material_id == ""):
				for r in range(r,(extra_data_row-2)):
					if(sheet_material.cell_value(r+2,9) == DNA_sampleID and sheet_material.cell_value(r+2,c) != "" and str(sheet_material.cell_value(r+2,c)) not in DNA_material_id):
						if(DNA_material_id == ""):
							DNA_material_id = str(sheet_material.cell_value(r+2,c))
						else:
							DNA_material_id = DNA_material_id + "," + str(sheet_material.cell_value(r+2,c))
						tumor_content_nr = sheet_material.cell_value(r+2,2)
						continue
					if(RNA_sampleID != "" and sheet_material.cell_value(r+2,9) == RNA_sampleID and sheet_material.cell_value(r+2,c) != "" and str(sheet_material.cell_value(r+2,c)) not in RNA_material_id):
						if(RNA_material_id == ""):
							RNA_material_id = str(sheet_material.cell_value(r+2,c))
						else:
							RNA_material_id = RNA_material_id + "," + str(sheet_material.cell_value(r+2,c))
						continue
			if(sheet_material.cell_value(r,c) == col_extraction_hospital and extraction_hospital == ""):
				for r in range(r,(library_pre_row-2)):
					if(sheet_material.cell_value(r+2,8) == DNA_sampleID):
						extraction_hospital = str(sheet_material.cell_value(r+2,c))
						break
			if(sheet_material.cell_value(r,c) == col_batch_nr and batch_nr == ""):
				for r in range(r,(nrows_material-2)):
					if(sheet_material.cell_value(r+2,0) == DNA_sampleID):
						batch_nr = str(sheet_material.cell_value(r+2,c))
	open_exl_material.release_resources()
	if(ipd_consent == "0.0"):
		ipd_consent = ""
	if(ipd_age == "" and ipd_birth_date != ""):
		ipd_age = "<1"
	inclusion_site_list = {'R': 'Radium', 'U': 'Ullevål', 'C': 'Riksen', 'A': 'Ahus', 'D': 'Drammen', 'B': 'Bærum', 'G': 'Gjøvik', 'I': 'Hamar', 'L': 'Lillehammer', 'T': 'Vestfold', 'K': 'Sørlandet', 'Q': 'Østfold', 'V': 'Telemark', 'Y': 'Lovisenberg', 'H': 'Haukeland', 'S': 'Stavanger', 'E': 'Fonna', 'F': 'Førde', 'O': 'St.Olavs', 'M': 'Nord-trøndelag', 'J': 'Møre og Romsdal', 'N': 'Nord Norge', 'P': 'Nordland'}
	if("IKKE IMPRESS" in ipd_consent):
		inclusion_site = ""
	else:
		try:
			site_letter_code = ipd_consent[-6]
			inclusion_site = inclusion_site_list.get(site_letter_code)
		except:
			inclusion_site = "Inclusion site"


def get_RNA_material_id(InPreD_clinical_data_file,RNA_sampleID):
	RNA_material_id_exist = False
	with open(InPreD_clinical_data_file, 'r', encoding="ISO-8859-1") as f:
		for l in f:
			if(RNA_sampleID == l.split('\t')[0]):
				RNA_material_id = l.split('\t')[8]
				RNA_material_id_exist = True
	f.close()
	if(RNA_material_id_exist == False):
		print("Warning: The "+ RNA_sampleID + " does not exist in the menta file! The report will be generated without RNA sample material ID!")
		RNA_material_id = ""
	return RNA_material_id


def update_ppt_template_data(inpred_node,ipd_no,ipd_gender,ipd_age,ipd_diagnosis_year,DNA_material_id,RNA_material_id,ipd_consent,requisition_hospital,ipd_clinical_diagnosis,tumor_type,sample_type,sample_material,pipline,tumor_content,ppt_template,output_ppt_file):
	if(ipd_gender != "" and ipd_gender != "X"):
		gender = ipd_gender[0]
	else:
		gender = ""
	if(ipd_age != "" and ipd_age != "-" and ipd_age != "XX" and ipd_age != "<1"):
		age = str(int(float(ipd_age)))
	else:
		age = ipd_age
	try:
		sample = sample_type + '\n' + sample_material
	except:
		sample = ""
	today_date = time.strftime("%d", time.localtime())
	today_month = time.strftime("%b", time.localtime())
	today_year = time.strftime("%Y", time.localtime())
	today = today_date + '\n' + today_month.upper() + '\n' + today_year
	ppt = Presentation(ppt_template)
	indexs = [1,3,4,5,6]
	for index in indexs:
		slide = ppt.slides[index]
		textbox1 = slide.shapes.add_textbox(Inches(3.75), Inches(0.11), Inches(1.33), Inches(0.50))
		tf1 = textbox1.text_frame
		tf1.paragraphs[0].text = ipd_no
		tf1.paragraphs[0].font.size = Pt(24)
		tf1.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		tf1.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox2 = slide.shapes.add_textbox(Inches(8.99), Inches(0.02), Inches(0.45), Inches(0.55))
		tf2 = textbox2.text_frame
		tf2.paragraphs[0].text = today
		tf2.paragraphs[0].font.size = Pt(9)
		tf2.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
		tf2.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
		textbox3 = slide.shapes.add_textbox(Inches(7.23), Inches(0.52), Inches(2.46), Inches(0.21))
		tf3 = textbox3.text_frame
		tf3.paragraphs[0].text = pipline
		tf3.paragraphs[0].font.size = Pt(7)
		tf3.paragraphs[0].font.color.rgb = RGBColor(64,64,64)
		textbox4 = slide.shapes.add_textbox(Inches(0.50), Inches(1.47), Inches(0.86), Inches(0.25))
		tf4 = textbox4.text_frame
		tf4.paragraphs[0].text = requisition_hospital
		tf4.paragraphs[0].font.size = Pt(10)
		tf4.paragraphs[0].alignment = PP_ALIGN.CENTER
		tf4.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		textbox5 = slide.shapes.add_textbox(Inches(0.71), Inches(1.84), Inches(0.86), Inches(0.50))
		tf5 = textbox5.text_frame
		tf5.paragraphs[0].text = sample
		tf5.paragraphs[0].font.size = Pt(8)
		tf5.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox6 = slide.shapes.add_textbox(Inches(0.81), Inches(2.65), Inches(0.63), Inches(0.33))
		tf6 = textbox6.text_frame
		tf6.paragraphs[0].text = tumor_content
		tf6.paragraphs[0].font.size = Pt(14)
		tf6.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox7 = slide.shapes.add_textbox(Inches(5.77), Inches(0.19), Inches(0.86), Inches(0.33))
		tf7 = textbox7.text_frame
		if(index == 1 or ipd_clinical_diagnosis == "-" or ipd_clinical_diagnosis == ""):
			tf7.paragraphs[0].text = str(tumor_type)
		else:
			tf7.paragraphs[0].text = ipd_clinical_diagnosis
		tf7.paragraphs[0].font.size = Pt(14)
		tf7.paragraphs[0].font.italic = True
		tf7.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		tf7.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox11 = slide.shapes.add_textbox(Inches(0.85), Inches(1.12), Inches(0.48), Inches(0.27))
		tf11 = textbox11.text_frame
		tf11.paragraphs[0].text = ipd_diagnosis_year
		tf11.paragraphs[0].font.size = Pt(10)
		tf11.paragraphs[0].alignment = PP_ALIGN.LEFT
		textbox12 = slide.shapes.add_textbox(Inches(0.61), Inches(0.35), Inches(1.02), Inches(0.33))
		tf12 = textbox12.text_frame
		tf12.paragraphs[0].text = inpred_node
		tf12.paragraphs[0].font.size = Pt(14)
		tf12.paragraphs[0].alignment = PP_ALIGN.CENTER
		tf12.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		if(index == 1):
			gender_age = ""
			ipd_material_id_index = ""
			ipd_consent_index = ""
		textbox8 = slide.shapes.add_textbox(Inches(0.69), Inches(0.79), Inches(0.87), Inches(0.40))
		tf8 = textbox8.text_frame
		tf8.paragraphs[0].text = gender_age
		tf8.paragraphs[0].font.size = Pt(18)
		tf8.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox9 = slide.shapes.add_textbox(Inches(0.73), Inches(2.25), Inches(0.70), Inches(0.26))
		tf9 = textbox9.text_frame
		tf9.paragraphs[0].text = ipd_material_id_index
		tf9.paragraphs[0].font.size = Pt(5)
		tf9.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox10 = slide.shapes.add_textbox(Inches(2.10), Inches(0.11), Inches(1.07), Inches(0.50))
		tf10 = textbox10.text_frame
		tf10.paragraphs[0].text = ipd_consent_index
		tf10.paragraphs[0].font.size = Pt(14)
		tf10.paragraphs[0].alignment = PP_ALIGN.CENTER
		tf10.paragraphs[0].font.italic = True
		tf10.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		gender_age = gender + '/' + age + 'y'
		if(RNA_material_id != ""):
			ipd_material_id_index = "DNA:" + DNA_material_id + "\nRNA:" + RNA_material_id
		else:
			ipd_material_id_index = "DNA:" + DNA_material_id
		ipd_consent_index = "Trial ID\n" + ipd_consent 

	ppt.save(output_ppt_file)


def insert_image_to_ppt(DNA_sampleID,DNA_normal_sampleID,RNA_sampleID,DNA_image_path,RNA_image_path,output_ppt_file):
	DNA_image = []
	RNA_image = []
	image_mark = "sample_QC_plot.png"
	ppt = Presentation(output_ppt_file)
	slide = ppt.slides[4]
	shapes = slide.shapes

	for file in os.listdir(DNA_image_path):
		if(DNA_sampleID in file and image_mark in file):
			image = os.path.join(DNA_image_path,file)
			DNA_image.append(image)
		if(DNA_normal_sampleID != "" and DNA_normal_sampleID in file and image_mark in file):
			image = os.path.join(DNA_image_path,file)
			DNA_image.append(image)

	if(RNA_image_path != ""):
		for file in os.listdir(RNA_image_path):
			if(RNA_sampleID in file and image_mark in file):
				image = os.path.join(RNA_image_path,file)
				RNA_image.append(image)
	left0 =  Inches(3.20)
	width0 = Inches(3.30)
	height0 = Inches(1.75)
	if(DNA_image != ''):
		top = Inches(1.55)
		d = 0
		for images in DNA_image:
			left = left0 + d * width0
			pic = slide.shapes.add_picture(images,left,top,width0,height0)
			d = d + 1	
	if(RNA_image != ''):
		top = Inches(3.44)
		r = 0
		for images in RNA_image:
			left = left0 + r * width0
			pic = slide.shapes.add_picture(images,left,top,width0,height0)
			r = r + 1
	ppt.save(output_ppt_file)


def insert_table_to_ppt(table_data_file,slide_n,table_name,left_h,top_h,width_h,left_t,top_t,width_t,height_t,font_size,table_header,output_ppt_file,if_print_rowNo):
	table_file = open(table_data_file)
	lines = table_file.readlines()
	first_line = lines[0]
	rows = len(lines)
	first_line_cells = first_line.split('\t')
	cols = len(first_line_cells) -1
	ppt = Presentation(output_ppt_file)
	try:
		slide = ppt.slides[slide_n-1]
	except:
		slide = ppt.slides.add_slide(ppt.slide_layouts[6])
	shapes = slide.shapes
	left = Inches(left_t)
	top = Inches(top_t)
	width = Inches(width_t)
	height = Inches(height_t)
	table = shapes.add_table(rows,cols,left,top,width,height).table
	table_rows = rows-1
	
	for c in range(cols):
		table.cell(0,c).text = table_header[c]
		table.cell(0,c).text_frame.paragraphs[0].font.size = Pt(font_size)
	row = 1
	for line in open(table_data_file):
		if(line != first_line):
			line_cells = line.split('\t')
			for j in range(len(line_cells) - 1):
				table.cell(row,j).text = str(line_cells[j])
				table.cell(row,j).text_frame.paragraphs[0].font.size = Pt(font_size)
			row += 1	
	textbox = slide.shapes.add_textbox(Inches(left_h),Inches(top_h),Inches(width_h),Inches(0.25))
	tf = textbox.text_frame
	if(if_print_rowNo == True):
		tf.paragraphs[0].text = table_name +" (N=" + str(table_rows) + ")"
	else:
		tf.paragraphs[0].text = table_name
	tf.paragraphs[0].font.size = Pt(8)
	tf.paragraphs[0].font.bold = True
	tf.paragraphs[0].alignment = PP_ALIGN.CENTER

	ppt.save(output_ppt_file)
	data_nrows = table_rows
	return data_nrows


def update_ppt_variant_summary_table(data_nrows,DNA_sampleID,RNA_sampleID,TMB_DRUP_nr,TMB_DRUP_str,DNA_variant_summary_file,RNA_variant_summary_file,output_file_preMTB_AppendixTable,output_table_file_filterResults_AllReporVariants_CodingRegion,output_ppt_file):
	DNA_summary_file = open(DNA_variant_summary_file)
	global str_TMB_DRUP
	global TMB_TSO500
	global MSI_TSO500
	for line in DNA_summary_file:
		if(line.startswith(DNA_sampleID)):
			if(line.split('\t')[1] == 'NA'):
				 TMB_illumina = "TMB = NA"
			else:
				TMB_illumina = "TMB = " + line.split('\t')[1]
				TMB_TSO500 = line.split('\t')[1]
			if(line.split('\t')[2] == 'NA'):
				MSI_illumina = "MSI = NA"
			else:
				MSI_illumina = "MSI = " + line.split('\t')[2]
				MSI_TSO500 = line.split('\t')[2]

			if(line.split()[3] != 'NA'):
			# X(Y/Z) 
			# Z<40: evaluation not reliable. Z>=40 && X=<10: MSI/Stable. Z>=40 && X>=20: MSI/Unstable. Z>=40 && 10<X<20: MSI/Likely unstable.
				X = float(line.split()[3])
				YZ = line.split()[4]
				Y_str = YZ.split('/')[0]
				Y = int(Y_str.split('(')[1])
				Z_str = YZ.split('/')[1]	
				Z = int(Z_str.split(')')[0])
				if(Z >= 40 and X >= 20):
					msi_text = "MS"
					stable_text = "Unstable"
				if(Z >= 40 and X <= 10):
					msi_text = "MS"
					stable_text = "Stable"
				if(Z >= 40 and 10 < X < 20):
					msi_text = "MS"
					stable_text = "Likely unstable"
				if(Z < 40):
					msi_text = "MS"
					stable_text = "Not reliable"
				msi_stable = str(Y) + " unstable out of " + str(Z)
			else:
				msi_text = "-"
				stable_text = "NA"
				msi_stable = "Not reliable"
			stable_text_long = stable_text + '\n' + msi_stable

	DNA_summary_file.close()
	
	splicing = "splicing: Not Assayed"
	fusion = "fusion: Not Assayed"	
	if(RNA_sampleID != ""):
		RNA_summary_file = open(RNA_variant_summary_file)
		for line in RNA_summary_file:
			if(line.startswith(RNA_sampleID)):
				splice_variants_str = line.split('\t')[4]
				if(splice_variants_str == 'NA'):
					splicing = "splicing: None reported"
				else:
					splicing = "splicing: " + splice_variants_str.split('(')[0]
					pattern_splicing = "\|(.*?)\("
					splicing_end = re.findall(pattern_splicing, splice_variants_str)
					if(splicing_end):
						for splice in splicing_end:
							splicing += ',' + splice
				gene_fusion_str = line.split('\t')[5]
				if(gene_fusion_str == 'NA\n'):
					fusion = "fusion: None reported"
				else:
					fusion = "fusion: " + gene_fusion_str.split('(')[0]
					pattern_fusion = "\|(.*?)\("
					fusion_end = re.findall(pattern_fusion, gene_fusion_str)
					if(fusion_end):
						for fus in fusion_end:
							fusion += ';' + fus
		RNA_summary_file.close()

	table_file_coding_region = open(output_table_file_filterResults_AllReporVariants_CodingRegion)
	appendix_nrows = len(table_file_coding_region.readlines()) - 1

	table_file_preMTBTable_Appendix = open(output_file_preMTB_AppendixTable)
	preMTB_appendix_nrows = len(table_file_preMTBTable_Appendix.readlines()) - 1
	if(TMB_DRUP_str == "-1"):
		str_TMB_DRUP = "-"
	else:
		effect_panel_size = float(TMB_DRUP_str.split('/')[1])
		if(effect_panel_size < 1.14):
			str_TMB_DRUP = "-"
		else:
			str_TMB_DRUP = str(TMB_DRUP_nr)

	ppt = Presentation(output_ppt_file)
	indexs = [1,5,6]
	for index in indexs:
		slide = ppt.slides[index]
		shapes = slide.shapes
		textbox1 = slide.shapes.add_textbox(Inches(5.76), Inches(1.60), Inches(0.41), Inches(0.27))
		tf1 = textbox1.text_frame
		tf1.paragraphs[0].text = msi_text
		tf1.paragraphs[0].font.size = Pt(10)
		tf1.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox2 = slide.shapes.add_textbox(Cm(15.84), Cm(4.18), Cm(2.26), Cm(0.64))
		tf2 = textbox2.text_frame
		tf2.paragraphs[0].text = stable_text_long
		tf2.paragraphs[0].font.size = Pt(7)
		tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox3 = slide.shapes.add_textbox(Inches(3.66), Inches(1.27), Inches(0.30), Inches(0.22))
		tf3 = textbox3.text_frame
		tf3.paragraphs[0].text = str(data_nrows)
		tf3.paragraphs[0].font.size = Pt(7)
		textbox4 = slide.shapes.add_textbox(Inches(3.66), Inches(1.14), Inches(0.30), Inches(0.22))
		tf4 = textbox4.text_frame
		tf4.paragraphs[0].text = str(appendix_nrows)
		tf4.paragraphs[0].font.size = Pt(7)
		tf4.paragraphs[0].alignment = PP_ALIGN.CENTER
		textbox5 = slide.shapes.add_textbox(Inches(3.66), Inches(1.02), Inches(0.30), Inches(0.22))
		tf5 = textbox5.text_frame
		tf5.paragraphs[0].text = str(preMTB_appendix_nrows)
		tf5.paragraphs[0].font.size = Pt(7)
		tf5.paragraphs[0].alignment = PP_ALIGN.CENTER
		if(TMB_DRUP_nr >= 0 and TMB_DRUP_nr <= 5 and str_TMB_DRUP != ""):
			roundshape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(7.07), Cm(3.90), Cm(0.58), Cm(0.58))
			roundshape.line.color.rgb = RGBColor(255,165,0)
			textbox5 = slide.shapes.add_textbox(Inches(2.74), Inches(1.54), Inches(0.32), Inches(0.21))
			tf5 = textbox5.text_frame
			tf5.paragraphs[0].text = str_TMB_DRUP
			tf5.paragraphs[0].font.size = Pt(7)
			tf5.paragraphs[0].font.bold = True
			tf5.paragraphs[0].alignment = PP_ALIGN.CENTER
			tf5.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		if((TMB_DRUP_nr > 5 and TMB_DRUP_nr <= 20) or str_TMB_DRUP == ""):
			roundshape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(8.27), Cm(3.90), Cm(0.58), Cm(0.58))
			roundshape.line.color.rgb = RGBColor(255,165,0)
			textbox5 = slide.shapes.add_textbox(Cm(8.23), Cm(3.90), Cm(0.58), Cm(0.60))
			tf5 = textbox5.text_frame
			tf5.paragraphs[0].text = str_TMB_DRUP
			tf5.paragraphs[0].font.size = Pt(7)
			tf5.paragraphs[0].font.bold = True
			tf5.paragraphs[0].alignment = PP_ALIGN.CENTER
			tf5.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		if(TMB_DRUP_nr > 20 and str_TMB_DRUP != ""):
			roundshape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(10.26), Cm(3.90), Cm(0.58), Cm(0.58))
			roundshape.line.color.rgb = RGBColor(255,165,0)
			textbox5 = slide.shapes.add_textbox(Cm(10.20), Cm(3.95), Cm(0.58), Cm(0.60))
			tf5 = textbox5.text_frame
			tf5.paragraphs[0].text = str_TMB_DRUP
			tf5.paragraphs[0].font.size = Pt(7)
			tf5.paragraphs[0].font.bold = True
			tf5.paragraphs[0].alignment = PP_ALIGN.CENTER
			tf5.paragraphs[0].font.color.rgb = RGBColor(250,250,250)
		textbox6 = slide.shapes.add_textbox(Inches(6.23), Inches(1.06), Inches(0.97), Inches(0.19))
		tf6 = textbox6.text_frame
		tf6.paragraphs[0].text = splicing
		tf6.paragraphs[0].font.size = Pt(7) 
		tf6.paragraphs[0].alignment = PP_ALIGN.LEFT
		textbox7 = slide.shapes.add_textbox(Inches(6.23), Inches(1.26), Inches(0.97), Inches(0.19))
		tf7 = textbox7.text_frame	
		tf7.paragraphs[0].text = fusion
		tf7.paragraphs[0].font.size = Pt(7)
		tf7.paragraphs[0].alignment = PP_ALIGN.LEFT
		if(index == 1):
			textbox8 = slide.shapes.add_textbox(Inches(5.14), Inches(2.63), Inches(0.53), Inches(0.25))
			tf8 = textbox8.text_frame
			tf8.paragraphs[0].text = TMB_DRUP_str
			tf8.paragraphs[0].font.size = Pt(8)
			tf8.paragraphs[0].alignment = PP_ALIGN.LEFT
			textbox9 = slide.shapes.add_textbox(Inches(3.95), Inches(3.70), Inches(0.88), Inches(0.25))
			tf9 = textbox9.text_frame
			tf9.paragraphs[0].text = TMB_illumina
			tf9.paragraphs[0].font.size = Pt(8)
			tf9.paragraphs[0].alignment = PP_ALIGN.LEFT
			textbox10 = slide.shapes.add_textbox(Inches(4.90), Inches(3.84), Inches(1.14), Inches(0.25))
			tf10 = textbox10.text_frame
			tf10.paragraphs[0].text = MSI_illumina
			tf10.paragraphs[0].font.size = Pt(8)
			tf10.paragraphs[0].alignment = PP_ALIGN.CENTER
		if(index == 5):
			textbox11 = slide.shapes.add_textbox(Inches(6.03), Inches(2.14), Inches(0.59), Inches(0.25))
			tf11 = textbox11.text_frame
			tf11.paragraphs[0].text = TMB_DRUP_str
			tf11.paragraphs[0].font.size = Pt(9)
			tf11.paragraphs[0].alignment = PP_ALIGN.LEFT
			textbox12 = slide.shapes.add_textbox(Inches(6.22), Inches(2.29), Inches(0.88), Inches(0.25))
			tf12 = textbox12.text_frame
			tf12.paragraphs[0].text = TMB_illumina
			tf12.paragraphs[0].font.size = Pt(9)
			tf12.paragraphs[0].alignment = PP_ALIGN.LEFT		
	ppt.save(output_ppt_file)
	return stable_text


def remisse_mail_writer(remisse_file,ipd_no,ipd_consent,DNA_normal_sampleID,RNA_sampleID,extraction_hospital,ipd_material_id,TMB_DRUP,stable_text,sample_material,sample_type,sample_list,pipline):
	from docx import Document
	from docx.shared import Pt
	from docx.shared import RGBColor as docRGBColor
	from docx.enum.text import WD_ALIGN_PARAGRAPH
	impress_id = ipd_consent
	sample_type = sample_type.replace("\n", "")
	doc = Document()
	doc.styles['Normal'].font.name = 'Calibri'
	doc.styles['Normal'].font.size = Pt(12)
	pg1 = doc.add_paragraph()
	if(DNA_normal_sampleID != ""):
		if_normal = "med"
	else:
		if_normal = "uten"
	text1 = pg1.add_run("Enhet for studierelatert diagnostikk, Avdeling for patologi, Klinikk for laboratoriemedisin, Oslo Universitetssykehus\nMAL for utsvar av TSO500, tumorprøve " + if_normal + " normal prøve (gDNA)")
	text1.bold = True
	pg1.alignment = WD_ALIGN_PARAGRAPH.CENTER
	pg2 = doc.add_paragraph()
	text2 = pg2.add_run("Merket grønt = fylles inn av molekylærbiolog og/eller patolog manuelt\n")
	text2.font.color.rgb = docRGBColor(0,176,80)
	pg2.add_run("##########################################################\n\n")
	pg2.add_run("PASIENT ID: " + ipd_no + "/" + impress_id + "\n\n")
	pg2.add_run("DIAGNOSE:\n")
	if(DNA_normal_sampleID != ""):
		if_normal_sampleID = ") og DNA ekstrahert fra blod:\n\n"
	else:
		if_normal_sampleID = "):\n\n"
	if(RNA_sampleID != ""):
		if_RNA_sampleID = " og RNA "
	else:
		if_RNA_sampleID = " "
	pg2.add_run("Utvidet genpanelanalyse, TSO500, DNA" + if_RNA_sampleID + "(ekstrahert ved " + extraction_hospital + " fra " + ipd_material_id + ": ")
	text3 = pg2.add_run("diagnose")
	text3.font.color.rgb = docRGBColor(0,176,80)
	pg2.add_run(if_normal_sampleID)
	if(TMB_DRUP >= 0 and TMB_DRUP <= 5):
		TMB_position = "lav"
	if(TMB_DRUP > 5 and TMB_DRUP <= 20):
		TMB_position = "intermediær"
	if(TMB_DRUP > 20):
		TMB_position = "høy"
	if(stable_text == "Unstable"):
		stable_text = "Ustabil"
	if(stable_text == "Stable"):
		stable_text = "Stabil"
	if(stable_text == "Not reliable"):
		stable_text = "Inkonklusiv"
	pg2.add_run("Tumor mutasjonsbyrde (TMB) estimat og kategori: " + str(TMB_DRUP) + " mut/Mb; " + TMB_position + "\n" + "Mikrosatellitt (MS) status: " + stable_text + "\nDNA-kopitallsendringer (estimert kopitall): ")	
	text4 = pg2.add_run("Ingen kopitall av sikker klinisk betydning\n")
	text4.font.color.rgb = docRGBColor(0,176,80)
	pg2.add_run("Genfusjoner: ")
	text5 = pg2.add_run("Ingen funn\n")
	text5.font.color.rgb = docRGBColor(0,176,80)
	pg2.add_run("Somatiske punktmutasjoner/insersjoner/delesjoner: ")
	text6 = pg2.add_run("Ingen funn av sikker klinisk betydning\n\n")
	text6.font.color.rgb = docRGBColor(0,176,80)
	pg2.add_run("Se vurdering og vedlegg. \n\n\nVURDERING:\n")
	text7 = pg2.add_run("Funnene er diskutert med overlege XXXX på Mol-MDT-møtet XX.XX.2022.\n\nEttersom det ikke er funn med klinisk betydning har man ikke kalt inn behandlende lege til Mol-MDT-møte. Ta kontakt dersom noe er uklart.\n\nSom ledd i IMPRESS-Norway-studien ble det utført FoundationOne CDx liquid test, med funn av XXXXXXXXXX.\n\nKimbane funn som skal følges opp? XXXXXXX.\n\nDet var dessverre ikke tilstrekkelig mengde og/eller kvalitet av DNA/RNA til at sekvenseringsanalysen kunne gjennomføres.\n\n")
	text7.font.color.rgb = docRGBColor(0,176,80)
	pg2.add_run("Kun funn med klinisk/diagnostisk betydning er rapportert, men se beskrivelse av TSO500-analysen nedenfor og den vedlagte Mol-MDT-rapporten for utfyllende informasjon om testresultatet.\n\n\nMOLEKYLÆRPATOLOGISK UNDERSØKELSE:\nBIOMATERIALET: " + ipd_material_id + "(" + sample_material + "; " + sample_type + ")\n\nTEST PANEL: TruSight Oncology 500 panel (Illumina)\n\n")
	if(extraction_hospital == "Enhet for studierelatert diagnostikk, OUS"):
		pg2.add_run("Materialet ekstrahert fra biomaterialet ved OUS sykehus ble kvalitetssikret før dypsekvensering med respektive protokoller og analyse pipelines.\n")
	else:
		pg2.add_run("Materialet mottatt ferdig isolert fra " + extraction_hospital + " sykehus ble kvalitetssikret før dypsekvensering med respektive protokoller og analysepipelines.\n")
	for i in range(len(sample_list)):
		pg2.add_run(str(sample_list[i]) + "\n")
	pg2.add_run("\nANALYSE AV DNA\nMengde DNA analysert: 50ng\n\nUtført dypsekvensering for deteksjon av punktmutasjoner, indeler og kopitallsendringer ved bruk av TruSight Oncology 500 panel (Illumina) som inkluderer 523 gener for DNA-analyser. Analysen inkluderer estimering av tumor mutasjonsbyrde (TMB) og mikrosatelitt (MS) status.\n\nProgramvare og analyseparametere: TruSight Oncology 500 Local App og in-house bioinformatisk pipeline for kvalitetssikring og variantfiltrering (" + pipline + "). Referansegenom GRCh37 ble brukt for kartlegging av sekvenser. Analysen er kjørt i \"Tumor ")
	if(DNA_normal_sampleID != ""):
		pg2.add_run("normal\"-innstilling for å ekskludere kimbanevarianter. ")
	else:
		pg2.add_run("only\"-innstilling og populasjonsdatabaser er benyttet for å ekskludere frekvente kimbanevarianter. ")
	pg2.add_run("Kvalitetssikring blir gjennomgått per enhet (chip) og per prøve sekvensert. For estimering av TMB benyttes antall ikke-synonyme mutasjoner detektert innenfor kodende DNA-områder delt på antall Mb sekvensert. Filter for varianter som inngår i TMB-beregningen er satt til minimum 5% variant allelfrekvens og minimum 50 sekvensfragmenter som dekker mutasjonssete. TMB-klassifiseringen er som følger: \"lav TMB\" for <5 mut/Mb, \"intermediær TMB\" for 5-20 mut/Mb, and \"høy TMB\" for >20 mut/Mb. TSO500-panelet analyserer 130 predefinerte MSI-seter for vurdering av MS-status. Et minimum av 40 slike seter må være analyserbare for å kunne pålitelig konkludere MS-status. Kopitallsendringer rapporteres som hovedregel kun ved kopitall >6.\n\nANALYSE AV RNA\nMengde RNA analysert: 40ng\n\nUtført dypsekvensering for deteksjon av fusjonsgener og spleisevarianter ved bruk av TruSight Oncology 500 panel (Illumina) som inkluderer 56 gener for RNA-analyser. Referansegenom GRCh37 ble brukt for kartlegging av sekvenser.\n\nProgramvare og analyseparametere: TruSight Oncology 500 Local App og in-house bioinformatisk pipeline for kvalitetssikring og variantfiltrering (" + pipline + "). Referansegenom GRCh37 ble brukt for kartlegging av sekvenser. En påvist fusjon må ha minst 3 unike sekvenser (reads) som støtter funnet.\n\nRESULTAT: Se diagnosefeltet og vedlagte Mol-MDT-rapport.\n\nResultater og tolkning er i henhold til, og innenfor rammene av kvalitet på prøvemateriale, det genomiske dekningsområdet til genpanelet, metodologi og anvendte kunnskapsdatabaser ved analysetidspunkt. Den operasjonelle pipelinen for TSO500-analyser ved InPreD OUS er i en utviklingsfase. Den kliniske signifikansen av TMB-verdien bør betraktes på bakgrunn av pasientens tumortype. En rapportering av funn med mulig terapeutisk implikasjon er ingen garanti eller lovnad om behandlingseffekt i pasienten. En klinisk helhetsvurdering av pasienten må foretas av behandlede lege ved mulig behandlingskonsekvens av analyseresultater. ")

	doc.save(remisse_file)


def update_clinical_master_file(InPreD_clinical_data_file,sample_id,if_generate_report,ipd_birth_year,clinical_diagnosis,ipd_gender,ipd_consent,material_id,ipd_collection_year,requisition_hospital,extraction_hospital,tumor_content_nr,batch_nr):
	global ipd_diagnosis_year
	global runID
	if_exist = False
	new_content = ""
	with open(InPreD_clinical_data_file, 'r', encoding="ISO-8859-1") as fr:
		for ln in fr:
			if(ln.split('\t')[0] == sample_id):
				if_exist = True
				line = sample_id + "\t" + runID + "\t" + if_generate_report + "\t" + ipd_birth_year + "\t" + ipd_diagnosis_year + "\t" + clinical_diagnosis + "\t" + ipd_gender[0] + "\t" + ipd_consent + "\t" + material_id + "\t" + ipd_collection_year + "\t" + requisition_hospital + "\t" + extraction_hospital + "\t" + str(tumor_content_nr) + "\t" +batch_nr + "\n"
				new_line = ln.replace(ln,line)
				new_content = new_content + new_line
			else:
				new_content = new_content + ln
	fr.close()
	if(if_exist == False):
		line = sample_id + "\t" + runID + "\t" + if_generate_report + "\t" + ipd_birth_year + "\t" + ipd_diagnosis_year + "\t" + clinical_diagnosis + "\t" + ipd_gender[0] + "\t" + ipd_consent + "\t" + material_id + "\t" + ipd_collection_year + "\t" + requisition_hospital + "\t" + extraction_hospital + "\t" + str(tumor_content_nr) + "\t" +batch_nr + "\n"
		with open(InPreD_clinical_data_file, 'a', encoding="ISO-8859-1") as fa:
			fa.write(line)
		fa.close()
	else:
		with open(InPreD_clinical_data_file, 'w', encoding="ISO-8859-1") as fw:
			fw.write(new_content)
		fw.close()


def update_clinical_tsoppi_file(InPreD_clinical_tsoppi_data_file,sample_id,if_generate_report,ipd_birth_year,clinical_diagnosis,ipd_gender,ipd_consent,material_id,ipd_collection_year,requisition_hospital,extraction_hospital,tumor_content_nr,batch_nr,sample_material,sample_type,tumor_type,TMB_DRUP,TMB_TSO500,MSI_TSO500,pipline):
	if_exist = False
	assay_name = ""
	nucleicacid = ""
	RNA_DNA_tumor_normal = ""
	global ipd_diagnosis_year
	global runID
	try:
		sample_type = sample_type.replace("\n", "")
	except:
		sample_type = ""
	if(pipline != "-" and pipline != ""):
		pipline = pipline.split(": ")[1]
	new_content = ""
	with open(InPreD_clinical_tsoppi_data_file, 'r') as fr:
		for ln in fr:
			if(ln.split('\t')[0] == sample_id):
				if_exist = True
				line = sample_id + "\t" + runID + "\t" + if_generate_report + "\t" + ipd_birth_year + "\t" + ipd_diagnosis_year + "\t" + clinical_diagnosis + "\t" + ipd_gender[0] + "\t" + ipd_consent + "\t" + material_id + "\t" + ipd_collection_year + "\t" + requisition_hospital + "\t" + extraction_hospital + "\t" + str(tumor_content_nr) + "\t" + batch_nr + "\t" + sample_material + "\t" + sample_type + "\t" + tumor_type + "\t" + str(TMB_DRUP) + "\t" + TMB_TSO500 + "\t" + MSI_TSO500 + "\t" + pipline + "\n"
				new_line = ln.replace(ln,line)
				new_content = new_content + new_line
			else:
				new_content = new_content + ln
	fr.close()
	if(if_exist == False):
		line = sample_id + "\t" + runID + "\t" + if_generate_report + "\t" + ipd_birth_year + "\t" + ipd_diagnosis_year + "\t" + clinical_diagnosis + "\t" + ipd_gender[0] + "\t" + ipd_consent + "\t" + material_id + "\t" + ipd_collection_year + "\t" + requisition_hospital + "\t" + extraction_hospital + "\t" + str(tumor_content_nr) + "\t" + batch_nr + "\t" + sample_material + "\t" + sample_type + "\t" + tumor_type + "\t" + str(TMB_DRUP) + "\t" + TMB_TSO500 + "\t" + MSI_TSO500 + "\t" + pipline + "\n"
		with open(InPreD_clinical_tsoppi_data_file, 'a') as fa:
			fa.write(line)
		fa.close()
	else:
		with open(InPreD_clinical_tsoppi_data_file, 'w') as fw:
			fw.write(new_content)
		fw.close()


def usage(exit_status = 0):
	print ("""Usage: python3  %s
        This script is a tool used to generate the paitent report based on the TSO500 analysis results and the personal intomation from the clinical data in In/InPreD_PRONTO_metadata.txt,
	and update the SOPPI results into the file Out/InPreD_PRONTO_metadata_tsoppi.txt when the reports are generated.
	This script could also fill the patient personal information into the clinical data file with the MTF files under the foder In/MTF/. (This fuction currently is only used by OUS)
	To run this script tool in your computer with python3, it will read the clinical data from In/InPreD_PRONTO_metadata.txt and generate reports for the Sample_id with Create_report==Y:
	
	python3 InPreD_PRONTO.py

	Extra parameters for OUS:
	-c, --clinical_file Fill the patient personal information into the clinical data file: InPreD_PRONTO_metadata.txt with the MTF files under the foder In/MTF/
	python3 InPreD_PRONTO.py -D <DNA_sampleID> -r <runID> -c
	or:
	python3 InPreD_PRONTO.py --DNAsampleID=<DNA_sampleID> --runID=<runID> --clinicalFile
	-m, --mail_draft Generate the Remisse_draft.docx file with report:
	python3 InPreD_PRONTO.py -m
	or:
	python3 InPreD_PRONTO.py --mailDraft

	This script will create sub-folder with runID/IPDXXX in Out/, move the IPD_Material file into it and generate all the results files under that sub-folder.
	 
	-h, --help See this help information and exit.
        """ % sys.argv[0])
	sys.exit(exit_status)


def main(argv):
	global runID
	global DNA_sampleID
	global RNA_sampleID
	global ipd_birth_year
	global ipd_clinical_diagnosis
	global ipd_gender
	global ipd_consent
	global ipd_material_id
	global DNA_material_id
	global RNA_material_id
	global ipd_collection_year
	global requisition_hospital
	global extraction_hospital
	global sample_material
	global sample_type
	global tumor_type
	global batch_nr
	global tumor_content_nr
	global TMB_DRUP
	global str_TMB_DRUP
	global TMB_TSO500
	global MSI_TSO500
	global pipline
	runID_RNA = ""
	runID_DNA = ""
	DNA_normal_sampleID = ""
	remisse_mail = False
	update_clinical_file = False
	try:
		opts, args = getopt.getopt(sys.argv[1:], "hr:D:mc", ["help", "runID=", "DNAsampleID=", "mailDraft", "clinicalFile"])
	except getopt.GetoptError:
		usage(1)

	for opt, arg in opts:
		if opt in ("-h", "--help"):
			usage()
		elif opt in ("-r", "--runID"):
			runID = arg
		elif opt in ("-D", "--DNAsampleID"):
			DNA_sampleID = arg
			ipd_no = DNA_sampleID.split('-')[0]
		elif opt in ("-m", "--mailDraft"):
			remisse_mail = True
		elif opt in ("-c", "--clinicalFile"):
			update_clinical_file = True
	runID_DNA = runID
	DNA_sampleID_format = "IP\w\d\d\d\d-D(\d|X)(\d|X)-\w(\d|X)(\d|X)-\w(\d|X)(\d|X)"

	base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
	config_file = base_dir + "/Config/configure_PRONTO.ini"
	InPreD_clinical_data_file = base_dir + "/In/InPreD_PRONTO_metadata.txt"
	output_path_root = base_dir + "/Out/"
	InPreD_clinical_tsoppi_data_file = base_dir + "/Out/InPreD_PRONTO_metadata_tsoppi.txt"
	cfg = ConfigParser()
	cfg.read(config_file)
	inpred_node = cfg.get("INPUT", "inpred_node")
	data_path = cfg.get("INPUT", "data_path")
	filter_col_nu_config = int(cfg.get("INPUT", "filter_col_nu"))

	if not os.path.exists(InPreD_clinical_data_file):
		print ("""      Error:
		The InPreD clinical file InPreD_PRONTO_metadata.txt does not exist!
		""")
		sys.exit(0)
	if(update_clinical_file == True):
		if not(re.fullmatch(DNA_sampleID_format, DNA_sampleID)):
			print("Warning: " + DNA_sampleID + " does not fit for the sample id format!")
		ipd_material_file = base_dir + "/In/MTF/" + ipd_no[:3] + '-' + ipd_no[3:] + "_Material Transit Form InPreD NGS.xlsx"
		if not os.path.exists(ipd_material_file):
			print ("""Error: IPD Material Transit Form InPreD NGS file does not exit under the MTF dir. PRONTO meta file could not be updated with patient personal information by parameter -c of this script!""")
			sys.exit(0)
		else:
			sample_list_file = data_path + runID_DNA + "_TSO_500_LocalApp_postprocessing_results/" + DNA_sampleID +"/" + "sample_list.tsv"
			for line in open(sample_list_file):
				if(line.startswith("RNA_tumor")):
					RNA_sampleID = line.split('\t')[1]
			get_patient_info_from_MTF(ipd_material_file,ipd_no,DNA_sampleID,RNA_sampleID)
			if_generate_report = "Y"
			update_clinical_master_file(InPreD_clinical_data_file,DNA_sampleID,if_generate_report,ipd_birth_year,ipd_clinical_diagnosis,ipd_gender,ipd_consent,DNA_material_id,ipd_collection_year,requisition_hospital,extraction_hospital,tumor_content_nr,batch_nr)
			print("Clinical data is added into PRONTO meta file for sample: " + DNA_sampleID)
			if(RNA_sampleID != ""):
				if_generate_report = "-"
				update_clinical_master_file(InPreD_clinical_data_file,RNA_sampleID,if_generate_report,ipd_birth_year,ipd_clinical_diagnosis,ipd_gender,ipd_consent,RNA_material_id,ipd_collection_year,requisition_hospital,extraction_hospital,tumor_content_nr,batch_nr)
				print("Clinical data is added into PRONTO meta file for sample: " + RNA_sampleID)
			sys.exit(0)
	ppt_nr = 0
	for ln in open(InPreD_clinical_data_file, encoding="ISO-8859-1"):
		if not(ln.startswith("#") or ln == ""):
			if(ln.split('\t')[2] == "Y"):
				ln = ln.replace("\n", "")
				global extra_path
				RNA_sampleID = ""
				sample_list = []
				DNA_sampleID = ln.split('\t')[0]
				runID_DNA = ln.split('\t')[1]
				runID = runID_DNA
				ipd_birth_year = ln.split('\t')[3]
				ipd_diagnosis_year = ln.split('\t')[4]
				ipd_clinical_diagnosis = ln.split('\t')[5]
				ipd_gender = ln.split('\t')[6]
				ipd_consent = ln.split('\t')[7]
				DNA_material_id = ln.split('\t')[8]
				ipd_collection_year = ln.split('\t')[9]
				requisition_hospital = ln.split('\t')[10]
				extraction_hospital = ln.split('\t')[11]
				tumor_content_nr = ln.split('\t')[12]
				batch_nr = ln.split('\t')[13]
				if not(re.fullmatch(DNA_sampleID_format, DNA_sampleID)):
					print("Warning: " + DNA_sampleID + " does not fit for the sample id format!")
				try:
					ipd_age = str(int(time.strftime("%Y", time.localtime())) - int(ipd_birth_year))
				except:
					ipd_age = "-"
				ipd_no = DNA_sampleID.split('-')[0]
				output_path = output_path_root + runID + "/" + DNA_sampleID + "/"
				extra_path = output_path + "extra_files"
				output_file_preMTB_table_path = output_path + DNA_sampleID
				today = time.strftime("%d %b %Y", time.localtime())
				sample_list_file = data_path + runID_DNA + "_TSO_500_LocalApp_postprocessing_results/" + DNA_sampleID +"/" + "sample_list.tsv"
				data_file_small_variant_table = data_path + runID_DNA + "_TSO_500_LocalApp_postprocessing_results/" + DNA_sampleID +"/" + DNA_sampleID + "_small_variant_table_forQC.tsv"
				if not os.path.exists(data_file_small_variant_table):
					data_file_small_variant_table = data_path + runID_DNA + "_TSO_500_LocalApp_postprocessing_results/" + DNA_sampleID +"/" + DNA_sampleID + "_small_variant_table.tsv"
				data_file_cnv_overview_plots = data_path + runID_DNA + "_TSO_500_LocalApp_postprocessing_results/" + DNA_sampleID +"/" + DNA_sampleID + "_CNV_overview_plots.pdf"

				if not os.path.exists(data_file_small_variant_table):
					print ("Error:  The data input file " + data_file_small_variant_table + " does not exist!")
					sys.exit(0)
				if not os.path.exists(output_path):
					os.makedirs(output_path)
				if not os.path.exists(extra_path):
					os.makedirs(extra_path)
				for line in open(sample_list_file):
					if not(line.startswith("#")):
						Sample_Type = line.split('\t')[0]
						Sample_ID = line.split('\t')[1]
						sample_list_str = Sample_Type + ": " + Sample_ID
						sample_list.append(sample_list_str)
					if(line.startswith("DNA_tumor")):
						DNA_run_dir = line.split('\t')[4]
					if(line.startswith("DNA_normal")):
						DNA_normal_sampleID = line.split('\t')[1]
					if(line.startswith("RNA_tumor")):
						RNA_run_dir = line.split('\t')[4]
						RNA_sampleID = line.split('\t')[1]		
					else:
						runID_RNA = ""
						RNA_sampleID = ""
						RNA_material_id = ""
				if(RNA_sampleID != ""):
					if (DNA_run_dir == RNA_run_dir):
						runID_RNA = runID_DNA
					else:
						RNA_run_dir_short = RNA_run_dir.split('/')[-1]
						runID_RNA = RNA_run_dir_short.split('_TSO')[0]
					RNA_material_id = get_RNA_material_id(InPreD_clinical_data_file,RNA_sampleID)
					ipd_material_id = "DNA:" + DNA_material_id + ",RNA:" + RNA_material_id
				else:
					ipd_material_id = "DNA:" + DNA_material_id

				for line in open(data_file_small_variant_table):
					if(line.startswith("#") and "Version string:" in line):
						pipline = line.split('\t')[0]
						pipline = pipline.split('] ')[1]
						pipline = pipline.split('\n')[0]
					if(line.startswith('#') and "Specified tumor purity" in line):
						if("not specified" in line or "a default value of 0.5 is being used" in line):
							tumor_content = "XX"
						else:
							tumor_content_float = Decimal(line.split(':')[1]).quantize(Decimal('0.00'))
							tumor_content = '~' + str(int(tumor_content_float*100)) + '%'
					if(line.startswith('#') and "Size of the target coding region" in line):
						target_cod_region = float(line.split(':')[1])
				MTB_format = False
				for i in range(0,filter_col_nu_config):
					j = str(i + 1)
					if(j == "2"):
						MTB_format = True
					filter_column = cfg.get("FILTER"+j, "filter_column")
					key_word = cfg.get("FILTER"+j, "key_word")
					all_col_output = cfg.get("FILTER"+j, "all_col_output")
					columns = cfg.get("FILTER"+j, "columns")
					output_table = cfg.get("FILTER"+j, "output_table")
					output_table_file_config_pre = output_file_preMTB_table_path + "_" + output_table + "_pre.txt"
					output_table_file_config = output_file_preMTB_table_path + "_" + output_table + ".txt"
					if(DNA_normal_sampleID != ""):
						columns = columns + ",AF_normal_DNA"
					if(all_col_output == "True"):
						if(j == "1"):
							filter1_min_depth_tumor_DNA = int(cfg.get("FILTER"+j, "min_depth_tumor_DNA"))
							all_data = read_exl(data_file_small_variant_table,filter_column,key_word)
							if(all_data == ""):
								all_data_config_DepthTumor_DNA = ""
							else:
								all_data_config_DepthTumor_DNA = filter_depth_tumor_all_col(all_data,filter1_min_depth_tumor_DNA)
							output_file_preMTB_workingTable_pre = output_file_preMTB_table_path + "_preMTB_workingTable_pre.txt"
							output_file_preMTB_workingTable = output_file_preMTB_table_path + "_preMTB_workingTable.txt"
							write_exl(output_file_preMTB_workingTable_pre,all_data_config_DepthTumor_DNA)
							clear_blank_line(output_file_preMTB_workingTable_pre,output_file_preMTB_workingTable)
						else:
							all_data = read_exl(output_file_preMTB_workingTable,filter_column,key_word)
							write_exl(output_table_file_config_pre,all_data)
							clear_blank_line(output_table_file_config_pre,output_table_file_config)
					if(j == "1"):
						filter1_min_depth_tumor_DNA = int(cfg.get("FILTER"+j, "min_depth_tumor_DNA"))
						data = read_exl_col(data_file_small_variant_table,filter_column,key_word,columns,MTB_format)
						if(data == ""):
							data_DepthTumor_DNA = ""
						else:
							data_DepthTumor_DNA = filter_depth_tumor_cols(data,filter1_min_depth_tumor_DNA)
						write_exl(output_table_file_config_pre,data_DepthTumor_DNA)
					else:
						data = read_exl_col(output_file_preMTB_workingTable,filter_column,key_word,columns,MTB_format)
						write_exl(output_table_file_config_pre,data)
					clear_blank_line(output_table_file_config_pre,output_table_file_config)
					MTB_format = False

				ppt_template = base_dir + "/In/Template/InPreD_MTB_template.pptx"
				DNA_variant_summary_file = data_path + runID_DNA + "_TSO_500_LocalApp_postprocessing_results/" + runID_DNA + "_variant_summary.tsv"
				if(runID_RNA != ""):
					RNA_variant_summary_file = data_path + runID_RNA + "_TSO_500_LocalApp_postprocessing_results/" + runID_RNA + "_variant_summary.tsv"
				else:
					RNA_variant_summary_file = ""
				output_ppt_file = output_path + DNA_sampleID + "_MTB_report.pptx"
				DNA_image_path = data_path + runID_DNA + "_TSO_500_LocalApp_postprocessing_results/" + DNA_sampleID +"/"
				if(RNA_sampleID != ""):
					RNA_image_path = DNA_image_path
				else:
					RNA_image_path = ""

				file = os.path.split(data_file_small_variant_table)[1]
				ipd_no_str = file.split('_')[0]
				try:
					sample_type_string = file.split('-')[2]
					sample_type_short = sample_type_string[0:1]
					sample_type_list = {'M': 'Metastasis', 'T': 'Primary Tumor', 'C': 'Cell-line', 'N': 'Normal/Control', 'P': 'Primary tumor\n naive', 'p': 'Primary tumor\n post-treatment', 'R': 'Regional met\n naive', 'r': 'Regional met\n post-treatment', 'D': 'Distal met\n naive', 'd': 'Distal met\n post-treatment', 'L': 'Liquid', 'E': 'naive', 'e': 'post treatment', 'A': 'post allo transplant', 'X': 'Unknown'}
					sample_type = sample_type_list.get(sample_type_short)
				except:
					sample_type = ""
				try:
					sample_material_string = file.split('-')[3]
					sample_material_short = sample_material_string[0:1]
					sample_material_list = {'F': 'Fresh Frozen', 'A': 'Archived FFPE', 'B': 'Blood', 'C': 'Cytology', 'M': 'Fresh bone marrow',  'E': 'Extramedullary','S': 'Buccal swab (normal)', 'X': 'Unspecified'}
					sample_material = sample_material_list.get(sample_material_short)
				except:
					sample_material = ""
				try:
					tumor_type_no = sample_material_string[1:3]
					tumor_type_list = {'00': 'Cancer origo incerta', '01': 'Adrenal Gland', '02': 'Ampulla of Vater', '03': 'Biliary Tract', '04': 'Bladder/Urinary Tract', '05': 'Bone', '06': 'Breast', '07': 'Cervix', '08': 'CNS/Brain', '09': 'Colon/Rectum', '10': 'Esophagus/Stomach', '11': 'Eye', '12': 'Head and Neck', '13': 'Kidney', '14': 'Liver', '15': 'Lung', '16': 'Lymphoid', '17': 'Myeloid', '18': 'Ovary/Fallopian Tube', '19': 'Pancreas', '20': 'Peripheral Nervous System', '21': 'Peritoneum', '22': 'Pleura', '23': 'Prostate', '24': 'Skin', '25': 'Soft Tissue', '26': 'Testis', '27': 'Thymus', '28': 'Thyroid', '29': 'Uterus', '30': 'Vulva/Vagina', 'XX': 'Not available'}
					tumor_type = tumor_type_list.get(tumor_type_no)
				except:
					tumor_type = ""
                
				min_AF_tumor_DNA = float(cfg.get("TMB", "min_AF_tumor_DNA"))
				min_depth_tumor_DNA = int(cfg.get("TMB", "min_depth_tumor_DNA"))
				TMB_filter_column = cfg.get("TMB", "TMB_filter_column")
				TMB_filter_key_word = cfg.get("TMB", "TMB_filter_key_word")
				TMB_coding_file_pre = output_path + DNA_sampleID + "_TMB_coding_pre.txt"
				TMB_coding_file = output_path + DNA_sampleID + "_TMB_coding.txt"
				TMB_DURP_coding_file_pre = output_path + DNA_sampleID + "_TMB_DURP_coding_pre.txt"
				TMB_DURP_coding_file = output_path + DNA_sampleID + "_TMB_DURP_coding.txt"
				TMB_DRUP_filter_key_word = cfg.get("TMB", "TMB_DRUP_filter_key_word")

				TMB_coding_data = read_exl(output_file_preMTB_workingTable,TMB_filter_column,TMB_filter_key_word)
				write_exl(TMB_coding_file_pre,TMB_coding_data)
				clear_blank_line(TMB_coding_file_pre,TMB_coding_file)
				TMB_DRUP_coding_data = read_exl(output_file_preMTB_workingTable,TMB_filter_column,TMB_DRUP_filter_key_word)
				write_exl(TMB_DURP_coding_file_pre,TMB_DRUP_coding_data)
				clear_blank_line(TMB_DURP_coding_file_pre,TMB_DURP_coding_file)

				rows_preMTB_AFTumor = 0
				for line in open(TMB_coding_file):
					line_cells = line.split('\t')
					if(line.startswith("Sample_ID")):
						for col in range(len(line_cells)):
							if(line_cells[col] == "AF_tumor_DNA"):
								col_AF_tumor = col
							if(line_cells[col] == "Depth_tumor_DNA"):
								col_Depth_tumor = col
					else:
						AF_tumor_DNA = float(line_cells[col_AF_tumor])
						Depth_tumor_DNA = int(line_cells[col_Depth_tumor])
						if(AF_tumor_DNA >= min_AF_tumor_DNA and Depth_tumor_DNA >= min_depth_tumor_DNA):
							rows_preMTB_AFTumor += 1
		
				if(target_cod_region == 0):
					TMB_In_house = -1
				else:
					TMB_In_house = round(rows_preMTB_AFTumor/target_cod_region, 1)

				rows_TMB_DRUP_AFTumor = 0
				for line in open(TMB_DURP_coding_file):
					line_cells = line.split('\t')
					if(line.startswith("Sample_ID")):
						for col in range(len(line_cells)):
							if(line_cells[col] == "AF_tumor_DNA"):
								col_DRUP_AF_tumor = col
							if(line_cells[col] == "Depth_tumor_DNA"):
								col_DRUP_Depth_tumor = col
					else:
						DRUP_AF_tumor_DNA = float(line_cells[col_DRUP_AF_tumor])
						DRUP_Depth_tumor_DNA = int(line_cells[col_DRUP_Depth_tumor])
						if(DRUP_AF_tumor_DNA >= min_AF_tumor_DNA and DRUP_Depth_tumor_DNA >= min_depth_tumor_DNA):
							rows_TMB_DRUP_AFTumor += 1

				if(target_cod_region == 0):
					TMB_DRUP = -1
					TMB_DRUP_str = "-1"
				else:
					TMB_DRUP = round(rows_TMB_DRUP_AFTumor/target_cod_region)
					TMB_DRUP_str = str(rows_TMB_DRUP_AFTumor) + '/' + str(target_cod_region)


				update_ppt_template_data(inpred_node,ipd_no,ipd_gender,ipd_age,ipd_diagnosis_year,DNA_material_id,RNA_material_id,ipd_consent,requisition_hospital,ipd_clinical_diagnosis,tumor_type,sample_type,sample_material,pipline,tumor_content,ppt_template,output_ppt_file)

				insert_image_to_ppt(DNA_sampleID,DNA_normal_sampleID,RNA_sampleID,DNA_image_path,RNA_image_path,output_ppt_file)

                		# Insert tables into PP file:
				slide8_table_header = ["Gene symbol", "MTB_format", "Ensembl_transcript_ID", "Genomic position", "Exon", "Protein change", "Change_summary", "Coding status", "Depth tumor DNA", "VAF tumor DNA", "Depth_normal DNA", "VAF normal DNA", "Illumina class", "Functional domain"]
				slide9_table_header = ["Gene symbol", "Genomic position", "DNA change", "cDNA change", "Protein change", "Change_summary", "Coding status", "Depth tumor DNA", "VAF tumor DNA", "Depth normal DNA", "VAF normal DNA", "Depth tumor RNA", "VAF tumor RNA", "Illumina class", "Class judgement comments"]
				if(DNA_normal_sampleID != ""):
					slide6_table_header = ["Gene symbol", "Protein change", "Coding status", "VAF tumor DNA [0,1]", "VAF normal DNA"]
				else:
					slide6_table_header = ["Gene symbol", "Protein change", "Coding status", "VAF tumor DNA [0,1]"]
        
                		# Slide2, slide6 and slide7 right side table: Variants that alter protein coding sequence
				slide6_table_data_file = output_file_preMTB_table_path + "_AllReporVariants_AltProtein.txt" 
				slide6_table_ppSlide = [2,6,7]
				slide6_table_name = "Variants that alter protein coding sequence "
				slide6_header_left = 7.36
				slide6_header_top = 0.82
				slide6_header_width = 2.55
				slide6_table_left = 7.23
				slide6_table_top = 1.06
				slide6_table_width = 2.76
				slide6_table_height = 1.63
				slide6_table_font_size = 7
				if_print_rowNo = False
				for table_index in slide6_table_ppSlide:
					slide6_table_nrows = insert_table_to_ppt(slide6_table_data_file,table_index,slide6_table_name,slide6_header_left,slide6_header_top,slide6_header_width,slide6_table_left,slide6_table_top,slide6_table_width,slide6_table_height,slide6_table_font_size,slide6_table_header,output_ppt_file,if_print_rowNo)
				output_file_preMTB_AppendixTable = output_file_preMTB_table_path + "_preMTBTable_Appendix.txt"
				output_table_file_filterResults_AllReporVariants_CodingRegion = output_file_preMTB_table_path + "_AllReporVariants_CodingRegion.txt"
				stable_text = update_ppt_variant_summary_table(slide6_table_nrows,DNA_sampleID,RNA_sampleID,TMB_DRUP,TMB_DRUP_str,DNA_variant_summary_file,RNA_variant_summary_file,output_file_preMTB_AppendixTable,output_table_file_filterResults_AllReporVariants_CodingRegion,output_ppt_file)
				output_file_preMTB_VigdisTable = output_file_preMTB_table_path + "_preMTBTable_Vigdis.txt"
                
				# Slide8 Table1: Sequence data summary: protein coding altering cariants*
				slide8_table_data_file = output_file_preMTB_VigdisTable
				slide8_table_ppSlide = 9
				slide8_table_name = "Table1: Sequence data summary: protein coding altering variants*  " + "  TSO500 " + pipline + " "
				slide8_header_left = 0.19
				slide8_header_top = 0.31
				slide8_header_width = 6.98
				slide8_table_left = 0.3
				slide8_table_top = 0.55
				slide8_table_width = 9.00
				slide8_table_height = 1.70
				slide8_table_font_size = 7
				if_print_rowNo = True
				slide8_table_nrows = insert_table_to_ppt(slide8_table_data_file,slide8_table_ppSlide,slide8_table_name,slide8_header_left,slide8_header_top,slide8_header_width,slide8_table_left,slide8_table_top,slide8_table_width,slide8_table_height,slide8_table_font_size,slide8_table_header,output_ppt_file,if_print_rowNo)

				# Slide9 Appendix 1: Sequence data summary: variants* called in sample IPDXXX
				slide9_table_data_file = output_file_preMTB_AppendixTable
				slide9_table_ppSlide = 10
				slide9_table_name = "Appendix 1: Sequence data summary: variants* called in sample " + DNA_sampleID
				slide9_header_left = 0.19
				slide9_header_top = 0.31
				slide9_header_width = 4.55
				slide9_table_left = 0.3
				slide9_table_top = 0.55
				slide9_table_width = 9.00
				slide9_table_height = 1.70
				slide9_table_font_size = 7
				if_print_rowNo = True
				slide9_table_nrows = insert_table_to_ppt(slide9_table_data_file,slide9_table_ppSlide,slide9_table_name,slide9_header_left,slide9_header_top,slide9_header_width,slide9_table_left,slide9_table_top,slide9_table_width,slide9_table_height,slide9_table_font_size,slide9_table_header,output_ppt_file,if_print_rowNo)

        			# Change slides order.
				ppt = Presentation(output_ppt_file)
				slides = ppt.slides._sldIdLst
				slides_list = list(slides)
				slides.remove(slides_list[7])
				slides.insert(9,slides_list[7])
				ppt.save(output_ppt_file)
				print("Generate report for " + DNA_sampleID)
				ppt_nr += 1

        			# Move TXT files generated by this script into extra_files folder.
				txt_files = os.listdir(output_path)
				for txt_file in txt_files:
					if(txt_file.endswith('.txt')):
						txt_file_path = os.path.join(output_path,txt_file)
						txt_file_extra_path = os.path.join(extra_path,txt_file)
						if os.path.exists(txt_file_extra_path):
							os.remove(txt_file_extra_path)
						shutil.move(txt_file_path, extra_path)

				# Move small variant data file and plots into the report folder.
				data_file_small_variant_table_cp = output_path + DNA_sampleID + "_small_variant_table_forQC.tsv"
				shutil.copyfile(data_file_small_variant_table, data_file_small_variant_table_cp)
				data_file_cnv_overview_plots_cp = output_path + DNA_sampleID + "_CNV_overview_plots.pdf"
				shutil.copyfile(data_file_cnv_overview_plots, data_file_cnv_overview_plots_cp)

				if(remisse_mail == True):
					remisse_file = output_path + ipd_no + "_Remisse_draft.docx"
					remisse_mail_writer(remisse_file,ipd_no,ipd_consent,DNA_normal_sampleID,RNA_sampleID,extraction_hospital,ipd_material_id,TMB_DRUP,stable_text,str(sample_material),sample_type,sample_list,pipline)
				ipd_material_file = base_dir + "/In/MTF/" + ipd_no[:3] + '-' + ipd_no[3:] + "_Material Transit Form InPreD NGS.xlsx"
				if os.path.exists(ipd_material_file):
					move_ipd_material_file = shutil.move(ipd_material_file, extra_path)
				if os.path.exists(InPreD_clinical_tsoppi_data_file):
					DNA_if_generate_report = "-"
					update_clinical_tsoppi_file(InPreD_clinical_tsoppi_data_file,DNA_sampleID,DNA_if_generate_report,ipd_birth_year,ipd_clinical_diagnosis,ipd_gender,ipd_consent,DNA_material_id,ipd_collection_year,requisition_hospital,extraction_hospital,tumor_content_nr,batch_nr,str(sample_material),sample_type,str(tumor_type),str_TMB_DRUP,TMB_TSO500,MSI_TSO500,pipline)
					if(RNA_sampleID != ""):
						RNA_if_generate_report = "-"
						RNA_str_TMB_DRUP = "-"
						RNA_TMB_TSO500 = "-"
						RNA_MSI_TSO500 = "-"
						RNA_pipline = "-"
						update_clinical_tsoppi_file(InPreD_clinical_tsoppi_data_file,RNA_sampleID,RNA_if_generate_report,ipd_birth_year,ipd_clinical_diagnosis,ipd_gender,ipd_consent,RNA_material_id,ipd_collection_year,requisition_hospital,extraction_hospital,tumor_content_nr,batch_nr,sample_material,sample_type,tumor_type,RNA_str_TMB_DRUP,RNA_TMB_TSO500,RNA_MSI_TSO500,RNA_pipline)
	if(ppt_nr > 1):	
		print("Go through the InPreD_PRONTO_metadata file, " + str(ppt_nr) +" reports are generated.")
	else:
		print("Go through the InPreD_PRONTO_metadata file, " + str(ppt_nr) +" report is generated.")

if __name__ == '__main__':
    main(sys.argv[1:])
